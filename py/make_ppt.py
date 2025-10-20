from pptx import Presentation
from pptx.util import Pt

prs = Presentation()

def add_slide(title, contents):
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = title
    tf = slide.placeholders[1].text_frame
    for c in contents:
        p = tf.add_paragraph()
        p.text = c
        p.font.size = Pt(20)

add_slide("Web地図はどのように動いているのか？", [
    "― HTMLとJavaScriptの基礎構成 ―",
    "屋外実習：スマホで地図＋GPSを体験",
    "次のステップ：その仕組みを理解する",
    "例題：daimaru.html"
])

add_slide("HTMLファイルの基本構成", [
    "<!DOCTYPE html> ～ </html> の枠組みで構成される",
    "<head>：タイトルや外部ライブラリの読み込み",
    "<body>：実際に表示される部分（#mapなど）",
    "<script>：地図を動かすJavaScriptを記述"
])

add_slide("Leafletによる地図描画の流れ", [
    "1️⃣ L.map() → 地図オブジェクトを作成",
    "2️⃣ setView([緯度, 経度], ズーム) → 初期位置を設定",
    "3️⃣ L.tileLayer() → 背景地図を指定",
    "4️⃣ .addTo(map) → 地図に追加",
    "➡ これだけでWeb地図が表示される"
])

add_slide("JavaScriptの役割", [
    "HTML = 構造, CSS = 見た目, JavaScript = 動き",
    "navigator.geolocation.getCurrentPosition(...) によりGPSを取得",
    "取得した座標を Leaflet の marker に反映",
    "➡ Webページがプログラムとして動く"
])

add_slide("イベントとUI制御", [
    "L.Control.extend() によりカスタムボタンを作成",
    "click イベントで locateOnce() を呼び出す",
    "➡ ボタン操作などのインタラクティブ機能を追加できる"
])

add_slide("Webサービスとしての仕組み", [
    "ブラウザ ⇄ JavaScript ⇄ 地図サーバ・GPS",
    "HTMLを読み込み → JSが外部サービスへアクセス",
    "地図タイルや位置情報を取得し動的に描画",
    "➡ Web地図 = ブラウザ上の小さなWebアプリ"
])

add_slide("まとめと応用", [
    "HTML：構造を定義",
    "JavaScript：動作・インタラクションを制御",
    "Web地図はブラウザだけで動作可能なGIS",
    "応用例：観察地点マッピング・現地調査データ共有"
])

prs.save("WebMap_HTML_JS_Basics.pptx")
print("✅ WebMap_HTML_JS_Basics.pptx を作成しました。")
